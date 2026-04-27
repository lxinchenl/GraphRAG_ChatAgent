from __future__ import annotations

import hashlib
import io
from pathlib import Path
from typing import Callable

import fitz
from docx import Document as DocxDocument
from PIL import Image
from pptx import Presentation

from .models import ParsedDocument

try:
    from rapidocr_onnxruntime import RapidOCR
except Exception:  # pragma: no cover - optional dependency runtime guard
    RapidOCR = None


SUPPORTED_SUFFIXES = {".pdf", ".docx", ".doc", ".pptx", ".png", ".jpg", ".jpeg", ".webp"}


class DocumentParser:
    def __init__(self, progress_callback: Callable[[str], None] | None = None) -> None:
        self.ocr_engine = RapidOCR() if RapidOCR else None
        self.progress_callback = progress_callback

    def parse_path(self, path: str | Path) -> list[ParsedDocument]:
        file_path = Path(path)
        suffix = file_path.suffix.lower()
        self._log(f"[parse] 开始处理文件: {file_path.name} ({suffix})")

        if suffix == ".pdf":
            return self._parse_pdf(file_path)
        if suffix == ".docx":
            return self._parse_docx(file_path)
        if suffix == ".doc":
            return self._parse_doc_with_word(file_path)
        if suffix == ".pptx":
            return self._parse_pptx(file_path)
        if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
            return self._parse_image(file_path)
        raise ValueError(f"Unsupported file type: {file_path.suffix}")

    def parse_directory(self, directory: str | Path) -> list[ParsedDocument]:
        root = Path(directory)
        results: list[ParsedDocument] = []
        files = [path for path in sorted(root.rglob("*")) if path.is_file() and path.suffix.lower() in SUPPORTED_SUFFIXES]
        self._log(f"[parse] 在目录中发现 {len(files)} 个可处理文件: {root}")
        for index, path in enumerate(files, start=1):
            self._log(f"[parse] 文件进度 {index}/{len(files)}: {path.name}")
            if path.is_file() and path.suffix.lower() in SUPPORTED_SUFFIXES:
                results.extend(self.parse_path(path))
        self._log(f"[parse] 文档解析完成，共生成 {len(results)} 条文档记录")
        return results

    def _parse_pdf(self, path: Path) -> list[ParsedDocument]:
        doc = fitz.open(path)
        records: list[ParsedDocument] = []
        doc_id = self._file_id(path)
        self._log(f"[parse][pdf] {path.name} 共 {len(doc)} 页")

        for page_index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if len(text) < 40:
                self._log(f"[parse][pdf] 第 {page_index} 页文本较少，尝试 OCR")
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                image_bytes = pix.tobytes("png")
                text = self._ocr_image_bytes(image_bytes)

            if not text.strip():
                self._log(f"[parse][pdf] 第 {page_index} 页无有效文本，跳过")
                continue

            records.append(
                ParsedDocument(
                    source_path=str(path),
                    doc_id=doc_id,
                    title=path.stem,
                    text=text.strip(),
                    modality="pdf_page",
                    page_number=page_index,
                )
            )
        self._log(f"[parse][pdf] {path.name} 解析完成，保留 {len(records)} 页")
        return records

    def _parse_docx(self, path: Path) -> list[ParsedDocument]:
        doc = DocxDocument(path)
        text = "\n".join([para.text.strip() for para in doc.paragraphs if para.text.strip()])
        self._log(f"[parse][docx] {path.name} 提取到 {len(text)} 个字符")
        return self._single_record(path, text, "docx")

    def _parse_doc_with_word(self, path: Path) -> list[ParsedDocument]:
        try:
            import win32com.client  # type: ignore
        except Exception as exc:  # pragma: no cover - only for Windows desktop
            raise RuntimeError(f"Reading .doc requires pywin32 and Word. File: {path}") from exc

        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        document = word.Documents.Open(str(path))
        try:
            text = document.Content.Text
        finally:
            document.Close(False)
            word.Quit()
        self._log(f"[parse][doc] {path.name} 提取到 {len(text)} 个字符")
        return self._single_record(path, text, "doc")

    def _parse_pptx(self, path: Path) -> list[ParsedDocument]:
        prs = Presentation(path)
        records: list[ParsedDocument] = []
        doc_id = self._file_id(path)
        self._log(f"[parse][pptx] {path.name} 共 {len(prs.slides)} 页")
        for slide_idx, slide in enumerate(prs.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text.strip())
            text = "\n".join([line for line in lines if line])
            if text:
                records.append(
                    ParsedDocument(
                        source_path=str(path),
                        doc_id=doc_id,
                        title=path.stem,
                        text=text,
                        modality="pptx_slide",
                        page_number=slide_idx,
                    )
                )
        self._log(f"[parse][pptx] {path.name} 解析完成，保留 {len(records)} 页")
        return records

    def _parse_image(self, path: Path) -> list[ParsedDocument]:
        with path.open("rb") as f:
            image_bytes = f.read()
        self._log(f"[parse][image] {path.name} 开始 OCR")
        text = self._ocr_image_bytes(image_bytes)
        self._log(f"[parse][image] {path.name} OCR 提取到 {len(text)} 个字符")
        return self._single_record(path, text, "image")

    def _single_record(self, path: Path, text: str, modality: str) -> list[ParsedDocument]:
        cleaned = text.strip()
        if not cleaned:
            return []
        return [
            ParsedDocument(
                source_path=str(path),
                doc_id=self._file_id(path),
                title=path.stem,
                text=cleaned,
                modality=modality,
            )
        ]

    def _ocr_image_bytes(self, image_bytes: bytes) -> str:
        if not self.ocr_engine:
            self._log("[parse][ocr] 未检测到 OCR 引擎，返回空文本")
            return ""
        image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        result, _ = self.ocr_engine(image)
        if not result:
            self._log("[parse][ocr] OCR 未识别到文本")
            return ""
        return "\n".join([item[1] for item in result if len(item) > 1 and item[1]])

    @staticmethod
    def _file_id(path: Path) -> str:
        digest = hashlib.md5(str(path.resolve()).encode("utf-8")).hexdigest()
        return f"doc_{digest}"

    def _log(self, message: str) -> None:
        if self.progress_callback:
            self.progress_callback(message)
